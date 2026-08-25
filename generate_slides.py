"""Generate 'How Fast Is Your Kernel, Really?' slide deck."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------------------------------------------------------------------------
# Palette & constants
# ---------------------------------------------------------------------------
BG_DARK = RGBColor(0xFF, 0xFF, 0xFF)
BG_ALT = RGBColor(0xF8, 0xF9, 0xFA)
WHITE = RGBColor(0x1A, 0x1A, 0x2E)
BLUE = RGBColor(0x0E, 0x6E, 0xB8)
GREEN = RGBColor(0x05, 0x7A, 0x55)
AMBER = RGBColor(0xB4, 0x6A, 0x00)
LIGHT_GRAY = RGBColor(0x55, 0x55, 0x55)
DIM_GRAY = RGBColor(0x77, 0x77, 0x77)
TABLE_ROW_DARK = RGBColor(0xF0, 0xF4, 0xF8)
TABLE_ROW_LIGHT = RGBColor(0xE8, 0xEE, 0xF4)
TABLE_HEADER_BG = RGBColor(0x0E, 0x6E, 0xB8)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT = "Calibri"
MONO_FONT = "Consolas"

ASSETS = "/home/ahpoddar/.cursor/projects/mnt-podman-storage-ahpoddar/assets"
OUT = "/mnt/podman_storage/ahpoddar/torch-profiler-workshop/slides.pptx"

TOTAL_SLIDES = 33

# Map screenshot filenames by rough purpose (DeepGEMM profiling)
if os.path.isdir(ASSETS):
    DEEPGEMM_SCREENSHOTS = sorted(
        [os.path.join(ASSETS, f) for f in os.listdir(ASSETS) if f.endswith(".png")]
    )
else:
    DEEPGEMM_SCREENSHOTS = []

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name=FONT, line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(font_size * 0.3)
    return tf


def add_paragraph(tf, text, font_size=18, color=WHITE, bold=False,
                  alignment=PP_ALIGN.LEFT, font_name=FONT, space_before=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if space_before:
        p.space_before = Pt(space_before)
    return p


def add_bullet_list(slide, left, top, width, height, items, font_size=18,
                    color=WHITE, bullet_color=BLUE, font_name=FONT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        p.level = 0
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        buClr = pPr.makeelement(qn("a:buClr"), {})
        srgbClr = buClr.makeelement(qn("a:srgbClr"), {"val": str(bullet_color)})
        buClr.append(srgbClr)
        pPr.append(buClr)
        buChar = pPr.makeelement(qn("a:buChar"), {"char": "\u2022"})
        pPr.append(buChar)
    return tf


def add_accent_bar(slide, left, top, width, height, color=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color, text="",
                     font_size=14, text_color=WHITE, border_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.name = FONT
    p.font.bold = True
    shape.text_frame.margin_left = Pt(8)
    shape.text_frame.margin_right = Pt(8)
    shape.text_frame.margin_top = Pt(6)
    shape.text_frame.margin_bottom = Pt(6)
    return shape


def add_table(slide, left, top, width, height, rows, cols, data,
              col_widths=None):
    """data is list-of-lists [row][col]. First row is header."""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r in range(rows):
        row_obj = table.rows[r]
        row_obj.height = Pt(36) if r == 0 else Pt(32)
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c]
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(14) if r == 0 else Pt(13)
            p.font.name = FONT
            p.font.bold = (r == 0)
            p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.text_frame.margin_left = Pt(6)
            cell.text_frame.margin_right = Pt(6)

            # Colors
            if r == 0:
                p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                cell_fill = cell.fill
                cell_fill.solid()
                cell_fill.fore_color.rgb = TABLE_HEADER_BG
            else:
                p.font.color.rgb = WHITE
                cell_fill = cell.fill
                cell_fill.solid()
                cell_fill.fore_color.rgb = TABLE_ROW_DARK if r % 2 == 1 else TABLE_ROW_LIGHT

            # Remove cell borders
            tcPr = cell._tc.get_or_add_tcPr()
            for border_name in ["a:lnL", "a:lnR", "a:lnT", "a:lnB"]:
                ln = tcPr.makeelement(qn(border_name), {"w": "0"})
                noFill = ln.makeelement(qn("a:noFill"), {})
                ln.append(noFill)
                tcPr.append(ln)

    return table_shape


def add_demo_placeholder(slide, text="[LIVE DEMO]"):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2), Inches(2.5), Inches(9.333), Inches(3.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
    shape.line.color.rgb = BLUE
    shape.line.width = Pt(2)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(32)
    p.font.color.rgb = BLUE
    p.font.bold = True
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_after = Pt(0)
    shape.text_frame.margin_top = Inches(1.2)
    return shape


def add_screenshot_placeholder(slide, text="[SCREENSHOT]", left=None,
                               top=None, width=None, height=None):
    """Gray rectangle with dashed-style border and centered placeholder text."""
    left = left or Inches(2.0)
    top = top or Inches(2.0)
    width = width or Inches(9.333)
    height = height or Inches(3.5)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE8, 0xEE, 0xF4)
    shape.line.color.rgb = DIM_GRAY
    shape.line.width = Pt(2)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(24)
    p.font.color.rgb = DIM_GRAY
    p.font.bold = True
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER
    return shape


def add_screenshot(slide, img_path, left, top, width, height=None):
    if os.path.exists(img_path):
        pic = slide.shapes.add_picture(img_path, left, top, width, height)
        return pic
    return None


def slide_title(slide, text, subtitle=None, top=Inches(0.3)):
    add_accent_bar(slide, Inches(0.6), top, Inches(0.12), Inches(0.55), BLUE)
    add_textbox(slide, Inches(0.9), top - Inches(0.05), Inches(11), Inches(0.7),
                text, font_size=32, color=BLUE, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.9), top + Inches(0.55), Inches(11), Inches(0.4),
                    subtitle, font_size=16, color=LIGHT_GRAY)


def slide_number(slide, num, total=TOTAL_SLIDES):
    add_textbox(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
                f"{num}/{total}", font_size=11, color=DIM_GRAY,
                alignment=PP_ALIGN.RIGHT)

# ---------------------------------------------------------------------------
# Build presentation
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]  # blank

slide_num = 0

# ===== SLIDE 1: Title =====
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)

add_accent_bar(sl, Inches(0), Inches(0), SLIDE_W, Inches(0.06), BLUE)

add_textbox(sl, Inches(1.5), Inches(1.3), Inches(10.3), Inches(1.2),
            "How Fast Is Your Kernel, Really?",
            font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(sl, Inches(1.5), Inches(2.6), Inches(10.3), Inches(0.8),
            "Profiling SOTA GPU Kernels with torch.profiler and Nsight Compute",
            font_size=22, color=BLUE, alignment=PP_ALIGN.CENTER)

add_accent_bar(sl, Inches(5.5), Inches(3.7), Inches(2.3), Inches(0.04), BLUE)

add_textbox(sl, Inches(1.5), Inches(3.9), Inches(10.3), Inches(0.6),
            "Aheli Poddar",
            font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(sl, Inches(1.5), Inches(4.4), Inches(10.3), Inches(0.6),
            "Associate Software Engineer, Red Hat PyTorch",
            font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(sl, Inches(1.5), Inches(5.0), Inches(10.3), Inches(0.6),
            "IISC Bangalore",
            font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

hw_box = add_rounded_rect(sl, Inches(4.0), Inches(5.8), Inches(5.3), Inches(0.7),
                          RGBColor(0xF0, 0xF4, 0xF8),
                          "NVIDIA H200  \u2022  SM 9.0  \u2022  143 GB HBM3e",
                          font_size=16, text_color=GREEN, border_color=GREEN)

add_accent_bar(sl, Inches(0), Inches(7.44), SLIDE_W, Inches(0.06), BLUE)
slide_number(sl, slide_num)


# ===== SLIDE 2: What We'll Cover =====
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
slide_title(sl, "What We\u2019ll Cover")

col_data = [
    ("DeepGEMM FP8", "3.09x over bf16 cuBLAS", "JIT-compiled FP8 GEMM\nfrom DeepSeek", BLUE),
    ("FlashAttention-3", "2.0\u20132.6x over FA-2", "Hopper-native attention\nwith WGMMA + TMA", GREEN),
    ("SonicMoE", "35 TFLOPS on H200", "IO-aware MoE kernels\nfrom Tri Dao", AMBER),
]

for i, (title, speedup, desc, accent) in enumerate(col_data):
    x = Inches(0.8 + i * 4.1)
    y = Inches(1.6)
    w = Inches(3.6)

    add_accent_bar(sl, x, y, w, Inches(0.05), accent)

    add_rounded_rect(sl, x, y + Inches(0.1), w, Inches(3.2),
                     RGBColor(0xF0, 0xF4, 0xF8), border_color=accent)

    add_textbox(sl, x + Inches(0.2), y + Inches(0.3), w - Inches(0.4), Inches(0.5),
                title, font_size=22, color=accent, bold=True, alignment=PP_ALIGN.CENTER)

    add_textbox(sl, x + Inches(0.2), y + Inches(0.9), w - Inches(0.4), Inches(0.5),
                speedup, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_textbox(sl, x + Inches(0.2), y + Inches(1.6), w - Inches(0.4), Inches(1.0),
                desc, font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

tf = add_textbox(sl, Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.6),
                 "Each kernel:  Code  \u2192  Perfetto Trace  \u2192  NCU Deep Dive",
                 font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

slide_number(sl, slide_num)


# ===== SLIDE 3: The Profiling Stack =====
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
slide_title(sl, "The Profiling Stack")

add_rounded_rect(sl, Inches(1.0), Inches(1.8), Inches(3.5), Inches(0.8),
                 RGBColor(0xF0, 0xF4, 0xF8), "torch.profiler", 20, BLUE, BLUE)

shape = sl.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(2.5), Inches(2.75), Inches(0.5), Inches(0.5))
shape.fill.solid()
shape.fill.fore_color.rgb = BLUE
shape.line.fill.background()

add_rounded_rect(sl, Inches(1.0), Inches(3.4), Inches(3.5), Inches(0.8),
                 RGBColor(0xF0, 0xF4, 0xF8), "Perfetto UI", 20, WHITE, BLUE)

add_textbox(sl, Inches(1.0), Inches(4.4), Inches(3.5), Inches(0.6),
            "Temporal view: WHEN kernels run",
            font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_rounded_rect(sl, Inches(8.8), Inches(1.8), Inches(3.5), Inches(0.8),
                 RGBColor(0xF0, 0xF4, 0xF8), "Nsight Compute (ncu)", 20, GREEN, GREEN)

shape = sl.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(10.3), Inches(2.75), Inches(0.5), Inches(0.5))
shape.fill.solid()
shape.fill.fore_color.rgb = GREEN
shape.line.fill.background()

add_rounded_rect(sl, Inches(8.8), Inches(3.4), Inches(3.5), Inches(0.8),
                 RGBColor(0xF0, 0xF4, 0xF8), "NCU GUI", 20, WHITE, GREEN)

add_textbox(sl, Inches(8.8), Inches(4.4), Inches(3.5), Inches(0.6),
            "Hardware view: HOW kernels use the GPU",
            font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_accent_bar(sl, Inches(2.5), Inches(5.6), Inches(8.3), Inches(0.04), AMBER)
add_textbox(sl, Inches(2.0), Inches(5.9), Inches(9.3), Inches(0.6),
            "torch.profiler tells you WHAT happened.  NCU tells you WHY.",
            font_size=22, color=AMBER, bold=True, alignment=PP_ALIGN.CENTER)

slide_number(sl, slide_num)


# ===== SLIDE 4: Kernel 1 — DeepGEMM Title =====
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_ALT)

add_textbox(sl, Inches(0.8), Inches(0.5), Inches(3), Inches(0.5),
            "KERNEL 1", font_size=16, color=BLUE, bold=True)

add_accent_bar(sl, Inches(0), Inches(0), Inches(0.12), SLIDE_H, BLUE)

add_textbox(sl, Inches(1.0), Inches(2.0), Inches(11.3), Inches(1.2),
            "DeepGEMM: JIT-Compiled FP8 GEMM",
            font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)

add_textbox(sl, Inches(1.0), Inches(3.5), Inches(11.3), Inches(0.5),
            "github.com/deepseek-ai/DeepGEMM",
            font_size=18, color=LIGHT_GRAY)

add_rounded_rect(sl, Inches(1.0), Inches(4.8), Inches(6.5), Inches(0.8),
                 RGBColor(0xF0, 0xF4, 0xF8),
                 "3.09x faster than bf16 cuBLAS on H200",
                 font_size=22, text_color=GREEN, border_color=GREEN)

slide_number(sl, slide_num)


# ===== SLIDE 5: What is FP8? =====
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
slide_title(sl, "What is FP8?")

add_textbox(sl, Inches(0.9), Inches(1.5), Inches(11), Inches(0.5),
            "E4M3 Format:  [S] [E E E E] [M M M]",
            font_size=26, color=WHITE, bold=True)

items = [
    "3 mantissa bits \u2014 optimized for forward pass / inference",
    "Block-scaled quantization: every 128 elements gets an FP32 scale factor",
    "Replaces the per-tensor scaling of classic INT8 quantization",
]
add_bullet_list(sl, Inches(0.9), Inches(2.4), Inches(11), Inches(2.2),
                items, font_size=18, color=WHITE)

add_rounded_rect(sl, Inches(2.5), Inches(5.0), Inches(8.3), Inches(1.0),
                 RGBColor(0xF0, 0xF4, 0xF8),
                 "2x Tensor Core throughput, half the data size",
                 font_size=24, text_color=GREEN, border_color=GREEN)

slide_number(sl, slide_num)


# ===== SLIDE 6: DeepGEMM Architecture =====
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
slide_title(sl, "DeepGEMM Architecture")

items = [
    "JIT compilation: first call compiles CUDA kernel, subsequent calls use cache",
    "generate_normal \u2192 FP8 casting + TMA-aligned scale factors",
    "transpose_fp32 helper: rearranges scale factors for TMA (padded stride 33 avoids bank conflicts)",
    "fp8_gemm_nt \u2192 sm90_fp8_gemm_1d2d_impl kernel",
]
add_bullet_list(sl, Inches(0.9), Inches(1.5), Inches(11), Inches(3.5),
                items, font_size=20, color=WHITE)

flow_items = ["generate_normal", "\u2192", "transpose_fp32", "\u2192", "sm90_fp8_gemm_1d2d_impl"]
x_start = Inches(1.0)
for i, item in enumerate(flow_items):
    if item == "\u2192":
        add_textbox(sl, x_start + Inches(i * 2.3), Inches(5.3), Inches(0.5), Inches(0.5),
                    "\u2192", font_size=28, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    else:
        add_rounded_rect(sl, x_start + Inches(i * 2.3), Inches(5.2), Inches(2.2), Inches(0.7),
                         RGBColor(0xF0, 0xF4, 0xF8), item, 14, WHITE, BLUE)

slide_number(sl, slide_num)


# ===== SLIDE 7: DeepGEMM — Perfetto Trace =====
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
slide_title(sl, "DeepGEMM \u2014 Perfetto Trace")

if len(DEEPGEMM_SCREENSHOTS) >= 1:
    add_screenshot(sl, DEEPGEMM_SCREENSHOTS[0],
                   Inches(0.8), Inches(1.3), Inches(11.7), Inches(4.5))
    add_textbox(sl, Inches(0.9), Inches(6.0), Inches(11.5), Inches(0.5),
                "CPU lane: fp8_deepgemm annotation \u2192 GPU lane: transpose_fp32 + sm90_fp8_gemm_1d2d_impl",
                font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
else:
    add_demo_placeholder(sl, "[LIVE DEMO \u2014 Perfetto UI]\n\nDeepGEMM FP8 Trace")

add_textbox(sl, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.5),
            "Look for: ProfilerStep#2/3/4, kernel names, CPU-GPU gap",
            font_size=14, color=AMBER, alignment=PP_ALIGN.CENTER)

slide_number(sl, slide_num)


# ===== SLIDE 8: DeepGEMM — Perfetto: bf16 Baseline =====
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
slide_title(sl, "DeepGEMM \u2014 Perfetto: bf16 Baseline")

add_screenshot_placeholder(sl, "[SCREENSHOT: Perfetto bf16 trace]",
                           Inches(1.5), Inches(1.5), Inches(10.3), Inches(3.2))

tf = add_textbox(sl, Inches(0.9), Inches(5.0), Inches(11.5), Inches(0.4),
                 "CPU: bf16_matmul \u2192 aten::matmul \u2192 aten::mm",
                 font_size=16, color=WHITE, font_name=MONO_FONT)
add_paragraph(tf, "GPU: nvjet_tst_256x128_64x4_1x2_h_bz_coopA_NNT",
              font_size=16, color=WHITE, font_name=MONO_FONT)
add_paragraph(tf, "3 ProfilerSteps \u00d7 ~300\u00b5s each",
              font_size=16, color=AMBER, bold=True)

slide_number(sl, slide_num)


# ===== SLIDE 9: DeepGEMM — Perfetto: FP8 Trace =====
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
slide_title(sl, "DeepGEMM \u2014 Perfetto: FP8 Trace")

add_screenshot_placeholder(sl, "[SCREENSHOT: Perfetto FP8 trace]",
                           Inches(1.5), Inches(1.5), Inches(10.3), Inches(3.2))

tf = add_textbox(sl, Inches(0.9), Inches(5.0), Inches(11.5), Inches(0.4),
                 "CPU: fp8_deepgemm annotation",
                 font_size=16, color=WHITE, font_name=MONO_FONT)
add_paragraph(tf, "GPU: transpose_fp32 (8\u00b5s) + sm90_fp8_gemm_1d2d_impl (~195\u00b5s)",
              font_size=16, color=WHITE, font_name=MONO_FONT)
add_paragraph(tf, "Note: transpose_fp32 has padded stride 33 to avoid SMEM bank conflicts",
              font_size=15, color=AMBER, bold=True)

slide_number(sl, slide_num)


# ===== SLIDE 10: DeepGEMM — Kernel Name Decode =====
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
slide_title(sl, "Reading Kernel Names")

# Left column: bf16 kernel name breakdown
add_textbox(sl, Inches(0.5), Inches(1.2), Inches(3.0), Inches(0.4),
            "bf16 cuBLAS", font_size=18, color=BLUE, bold=True)

bf16_lines = [
    "nvjet_tst_256x128_64x4_1x2_h_bz_coopA_NNT",
    "",
    "nvjet       = NVIDIA JetBLAS (Hopper cuBLAS)",
    "tst         = tensor-sparse-tensor variant",
    "256x128     = CTA tile (rows \u00d7 cols)",
    "64x4        = 64 K-tile \u00d7 4 pipeline stages",
    "1x2         = warp arrangement",
    "h           = Hopper architecture",
    "bz          = block-zero optimization",
    "coopA       = cooperative launch on A",
    "NNT         = layout (A=row, B=row, out=col)",
]
tf = add_textbox(sl, Inches(0.4), Inches(1.7), Inches(6.3), Inches(5.5),
                 bf16_lines[0], font_size=12, color=WHITE, font_name=MONO_FONT)
for line in bf16_lines[1:]:
    add_paragraph(tf, line, font_size=12, color=WHITE if "=" not in line else LIGHT_GRAY,
                  font_name=MONO_FONT)

# Right column: FP8 kernel name breakdown
add_textbox(sl, Inches(7.0), Inches(1.2), Inches(3.0), Inches(0.4),
            "FP8 DeepGEMM", font_size=18, color=GREEN, bold=True)

fp8_lines = [
    "sm90_fp8_gemm_1d2d_impl",
    "",
    "sm90  = Hopper SM 9.0",
    "fp8   = FP8 precision",
    "gemm  = General Matrix Multiply",
    "1d    = 1D scale on A (per-token)",
    "2d    = 2D scale on B (per-128\u00d7128 block)",
    "impl  = implementation kernel",
]
tf = add_textbox(sl, Inches(7.0), Inches(1.7), Inches(5.8), Inches(5.5),
                 fp8_lines[0], font_size=13, color=WHITE, font_name=MONO_FONT)
for line in fp8_lines[1:]:
    add_paragraph(tf, line, font_size=13, color=WHITE if "=" not in line else LIGHT_GRAY,
                  font_name=MONO_FONT)

# Vertical divider
add_accent_bar(sl, Inches(6.65), Inches(1.5), Inches(0.03), Inches(5.0), DIM_GRAY)

slide_number(sl, slide_num)


# ===== SLIDE 11: DeepGEMM — NCU Comparison =====
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
slide_title(sl, "DeepGEMM \u2014 NCU Comparison")

data = [
    ["Metric", "bf16 cuBLAS", "DeepGEMM FP8"],
    ["Duration", "375 \u00b5s", "251 \u00b5s"],
    ["Compute Throughput", "91.6%", "67.1%"],
    ["Memory Throughput", "70.7%", "54.5%"],
    ["TMA Instructions", "344K", "59K"],
    ["Memory BW", "969 GB/s", "543 GB/s"],
    ["Wallclock (profiler)", "~600 \u00b5s", "~195 \u00b5s"],
]

add_table(sl, Inches(1.5), Inches(1.5), Inches(10.3), Inches(3.5),
          len(data), 3, data,
          col_widths=[Inches(4.0), Inches(3.15), Inches(3.15)])

add_textbox(sl, Inches(1.0), Inches(5.5), Inches(11.3), Inches(1.0),
            "bf16 is more efficient.  FP8 is faster.\nWhy?  2x Tensor Core FLOPS + half the data.",
            font_size=20, color=AMBER, bold=True, alignment=PP_ALIGN.CENTER)

slide_number(sl, slide_num)


# ===== SLIDE 12: NCU: bf16 Speed of Light =====
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
slide_title(sl, "NCU: bf16 GEMM \u2014 Speed of Light")

add_screenshot_placeholder(sl, "[SCREENSHOT: NCU bf16 Speed of Light chart]",
                           Inches(1.5), Inches(1.5), Inches(10.3), Inches(3.5))

add_textbox(sl, Inches(1.0), Inches(5.3), Inches(11.3), Inches(0.5),
            "91.6% Compute Throughput \u2014 this is what \u2018good\u2019 looks like",
            font_size=20, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(sl, Inches(1.0), Inches(5.9), Inches(11.3), Inches(0.5),
            "NCU flags: \u201cHigh Throughput \u2014 utilizing greater than 80%\u201d",
            font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

slide_number(sl, slide_num)


# ===== SLIDE 13: NCU: FP8 Speed of Light =====
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
slide_title(sl, "NCU: DeepGEMM FP8 \u2014 Speed of Light")

add_screenshot_placeholder(sl, "[SCREENSHOT: NCU FP8 Speed of Light chart]",
                           Inches(1.5), Inches(1.5), Inches(10.3), Inches(3.5))

add_textbox(sl, Inches(1.0), Inches(5.3), Inches(11.3), Inches(0.5),
            "67.1% Compute, 54.5% Memory \u2014 compute-bound with scale factor overhead",
            font_size=20, color=AMBER, bold=True, alignment=PP_ALIGN.CENTER)

slide_number(sl, slide_num)


# ===== SLIDE 14: NCU: Memory Architecture =====
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
slide_title(sl, "NCU: Memory Data Path Comparison")

# Left column: bf16
add_textbox(sl, Inches(0.6), Inches(1.2), Inches(2.0), Inches(0.4),
            "bf16", font_size=20, color=BLUE, bold=True)

bf16_mem_items = [
    "TMA: 344K instructions",
    "Data: 2.82 GB through pipeline",
    "DSMEM: 382 instructions (uses SM clusters)",
    "Memory BW: 969 GB/s",
]
add_bullet_list(sl, Inches(0.6), Inches(1.7), Inches(5.8), Inches(2.0),
                bf16_mem_items, font_size=16, color=WHITE, bullet_color=BLUE)

add_screenshot_placeholder(sl, "[SCREENSHOT: NCU bf16 memory diagram]",
                           Inches(0.6), Inches(4.0), Inches(5.8), Inches(2.8))

# Right column: FP8
add_textbox(sl, Inches(7.0), Inches(1.2), Inches(3.0), Inches(0.4),
            "FP8 DeepGEMM", font_size=20, color=GREEN, bold=True)

fp8_mem_items = [
    "TMA: 59K instructions (6x fewer)",
    "Data: 1.44 GB (half the data)",
    "DSMEM: 0 (no cross-SM sharing)",
    "Memory BW: 543 GB/s",
]
add_bullet_list(sl, Inches(7.0), Inches(1.7), Inches(5.8), Inches(2.0),
                fp8_mem_items, font_size=16, color=WHITE, bullet_color=GREEN)

add_screenshot_placeholder(sl, "[SCREENSHOT: NCU FP8 memory diagram]",
                           Inches(7.0), Inches(4.0), Inches(5.8), Inches(2.8))

# Vertical divider
add_accent_bar(sl, Inches(6.6), Inches(1.2), Inches(0.03), Inches(5.8), DIM_GRAY)

slide_number(sl, slide_num)


# ===== SLIDE 15: NCU: Occupancy =====
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
slide_title(sl, "NCU: Occupancy \u2014 Both Kernels Are Resource-Heavy")

data = [
    ["Metric", "bf16", "FP8"],
    ["Registers/Thread", "168", "168"],
    ["Dynamic SMEM/Block", "213 KB", "216 KB"],
    ["Theoretical Occupancy", "18.75%", "18.75%"],
    ["Achieved Occupancy", "14.72%", "14.00%"],
    ["Block Limit (Registers)", "1", "1"],
    ["Block Limit (Shared Mem)", "1", "1"],
    ["Grid Size", "132 (= 132 SMs)", "132"],
    ["Cluster Size", "2", "2"],
]

add_table(sl, Inches(1.5), Inches(1.3), Inches(10.3), Inches(4.5),
          len(data), 3, data,
          col_widths=[Inches(4.5), Inches(2.9), Inches(2.9)])

add_textbox(sl, Inches(1.0), Inches(6.1), Inches(11.3), Inches(0.8),
            "Both kernels use maximum on-chip resources. 1 block per SM. Low occupancy by design.",
            font_size=18, color=AMBER, bold=True, alignment=PP_ALIGN.CENTER)

slide_number(sl, slide_num)


# ===== SLIDE 16: DeepGEMM — NCU Key Findings =====
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
slide_title(sl, "DeepGEMM \u2014 NCU Key Findings")

findings = [
    ("Barrier Stalls", "4.3 cycles/warp (32.95% est. speedup)", BLUE),
    ("Occupancy", "14% achieved (18.75% theoretical) \u2014 limited by 168 regs + 216 KB SMEM", GREEN),
    ("Shared Store Bank Conflicts", "FP8: 12.55% vs bf16: 8.65% \u2014 FP8 pays extra for scale factor stores", AMBER),
]

for i, (label, detail, accent) in enumerate(findings):
    y = Inches(1.6 + i * 1.5)
    add_accent_bar(sl, Inches(0.9), y, Inches(0.08), Inches(1.0), accent)
    add_textbox(sl, Inches(1.2), y, Inches(10.5), Inches(0.5),
                label, font_size=22, color=accent, bold=True)
    add_textbox(sl, Inches(1.2), y + Inches(0.5), Inches(10.5), Inches(0.5),
                detail, font_size=17, color=WHITE)

add_textbox(sl, Inches(1.5), Inches(6.0), Inches(10.3), Inches(0.6),
            '"The scale factor overhead is the price of FP8 quantization"',
            font_size=18, color=AMBER, bold=True, alignment=PP_ALIGN.CENTER)

slide_number(sl, slide_num)


# ===== SLIDE 17: Kernel 2 — FlashAttention-3 Title =====
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_ALT)

add_textbox(sl, Inches(0.8), Inches(0.5), Inches(3), Inches(0.5),
            "KERNEL 2", font_size=16, color=GREEN, bold=True)

add_accent_bar(sl, Inches(0), Inches(0), Inches(0.12), SLIDE_H, GREEN)

add_textbox(sl, Inches(1.0), Inches(2.0), Inches(11.3), Inches(1.2),
            "FlashAttention-3: Hopper-Native Attention",
            font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)

add_textbox(sl, Inches(1.0), Inches(3.5), Inches(11.3), Inches(0.5),
            "github.com/Dao-AILab/flash-attention  (hopper/)",
            font_size=18, color=LIGHT_GRAY)

add_rounded_rect(sl, Inches(1.0), Inches(4.8), Inches(8.0), Inches(0.8),
                 RGBColor(0xF0, 0xF4, 0xF8),
                 "2.0\u20132.6x faster than FA-2 across all sequence lengths",
                 font_size=22, text_color=GREEN, border_color=GREEN)

slide_number(sl, slide_num)


# ===== SLIDE 18: FA-2 to FA-3 =====
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
slide_title(sl, "What Changed from FA-2 to FA-3")

features = [
    ("WGMMA", "Warp Group MMA: larger Tensor Core tiles", GREEN),
    ("TMA", "Tensor Memory Accelerator: async bulk HBM \u2192 SMEM", BLUE),
    ("Warp Specialization", "Producer warps (TMA) overlap with consumer warps (WGMMA)", AMBER),
]

for i, (name, desc, accent) in enumerate(features):
    y = Inches(1.6 + i * 1.3)
    num_box = add_rounded_rect(sl, Inches(0.9), y, Inches(0.6), Inches(0.6),
                               accent, str(i + 1), 20, WHITE)
    add_textbox(sl, Inches(1.8), y, Inches(3.5), Inches(0.5),
                name, font_size=22, color=accent, bold=True)
    add_textbox(sl, Inches(1.8), y + Inches(0.45), Inches(10), Inches(0.5),
                desc, font_size=16, color=WHITE)

add_textbox(sl, Inches(0.9), Inches(5.4), Inches(11.5), Inches(0.5),
            "Plus: Ping-pong scheduling overlaps GEMM and softmax",
            font_size=18, color=WHITE, bold=True)

slide_number(sl, slide_num)


# ===== SLIDE 19: FA-3 Scaling Results =====
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
slide_title(sl, "FA-3 vs FA-2 \u2014 Scaling Results")

data = [
    ["Seq Length", "FA-2 (ms)", "FA-3 (ms)", "Speedup"],
    ["512", "0.059", "0.031", "1.89x"],
    ["2048", "0.980", "0.474", "2.07x"],
    ["4096", "3.698", "1.667", "2.22x"],
    ["8192", "14.054", "6.593", "2.13x"],
]

add_table(sl, Inches(2.0), Inches(1.5), Inches(9.3), Inches(3.0),
          len(data), 4, data,
          col_widths=[Inches(2.3), Inches(2.3), Inches(2.3), Inches(2.4)])

add_textbox(sl, Inches(1.5), Inches(5.3), Inches(10.3), Inches(0.6),
            "Consistent 2x across all sequence lengths \u2014 not a one-off optimization",
            font_size=20, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

slide_number(sl, slide_num)


# ===== SLIDE 20: FA-3 Kernel Names =====
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
slide_title(sl, "FA-3 \u2014 Kernel Names Tell the Story")

add_textbox(sl, Inches(0.9), Inches(1.8), Inches(3), Inches(0.4),
            "FlashAttention-2", font_size=18, color=LIGHT_GRAY, bold=True)
add_rounded_rect(sl, Inches(0.9), Inches(2.3), Inches(5.5), Inches(1.0),
                 RGBColor(0xF0, 0xF4, 0xF8), "flash_fwd_kernel",
                 18, WHITE, BLUE)
add_textbox(sl, Inches(0.9), Inches(3.5), Inches(5.5), Inches(0.5),
            "Hand-written CUDA kernel, vendored in PyTorch",
            font_size=14, color=LIGHT_GRAY)

add_textbox(sl, Inches(7.0), Inches(1.8), Inches(3), Inches(0.4),
            "FlashAttention-3", font_size=18, color=LIGHT_GRAY, bold=True)
add_rounded_rect(sl, Inches(7.0), Inches(2.3), Inches(5.5), Inches(1.0),
                 RGBColor(0xF0, 0xF4, 0xF8), "device_kernel",
                 18, WHITE, GREEN)
add_textbox(sl, Inches(7.0), Inches(3.5), Inches(5.5), Inches(0.5),
            "CUTLASS 3.x CuTeDSL template, Hopper-native",
            font_size=14, color=LIGHT_GRAY)

add_textbox(sl, Inches(6.0), Inches(2.5), Inches(1.2), Inches(0.7),
            "\u2192", font_size=36, color=AMBER, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(sl, Inches(1.5), Inches(5.2), Inches(10.3), Inches(0.8),
            '"The name change reflects the architecture shift:\nfrom hand-tuned to template-generated"',
            font_size=20, color=AMBER, bold=True, alignment=PP_ALIGN.CENTER)

slide_number(sl, slide_num)


# ===== SLIDE 21: FA-3 Perfetto Trace =====
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
slide_title(sl, "FA-3 \u2014 Perfetto Trace")

add_demo_placeholder(sl, "[LIVE DEMO \u2014 Perfetto UI]\n\nFA-2 trace (flash_fwd_kernel)  vs  FA-3 trace (device_kernel)\nSame attention, different kernel, 2x faster")

slide_number(sl, slide_num)


# ===== SLIDE 22: FA-3 — NCU Comparison (NEW) =====
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
slide_title(sl, "FA-3 \u2014 NCU Comparison")

add_screenshot_placeholder(sl, "[SCREENSHOT: NCU FA-2 vs FA-3]",
                           Inches(1.5), Inches(1.5), Inches(10.3), Inches(4.5))

add_textbox(sl, Inches(1.0), Inches(6.3), Inches(11.3), Inches(0.5),
            "Compare: occupancy, warp stalls, memory throughput between FA-2 and FA-3 kernels",
            font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

slide_number(sl, slide_num)


# ===== SLIDE 23: Kernel 3 — SonicMoE Title =====
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_ALT)

add_textbox(sl, Inches(0.8), Inches(0.5), Inches(3), Inches(0.5),
            "KERNEL 3", font_size=16, color=AMBER, bold=True)

add_accent_bar(sl, Inches(0), Inches(0), Inches(0.12), SLIDE_H, AMBER)

add_textbox(sl, Inches(1.0), Inches(2.0), Inches(11.3), Inches(1.2),
            "SonicMoE: IO-Aware MoE Kernels",
            font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)

add_textbox(sl, Inches(1.0), Inches(3.5), Inches(11.3), Inches(0.5),
            "github.com/Dao-AILab/sonic-moe  (Tri Dao, Princeton / Together AI)",
            font_size=18, color=LIGHT_GRAY)

add_rounded_rect(sl, Inches(1.0), Inches(4.8), Inches(8.0), Inches(0.8),
                 RGBColor(0xF0, 0xF4, 0xF8),
                 "45% less activation memory, 1.86x faster than ScatterMoE",
                 font_size=22, text_color=AMBER, border_color=AMBER)

slide_number(sl, slide_num)


# ===== SLIDE 24: Why MoE Kernels Are Hard =====
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
slide_title(sl, "Why MoE Kernels Are Hard")

items = [
    "Each expert processes a variable number of tokens (varlen-M Grouped GEMM)",
    "Fine-grained MoEs: small expert size \u2192 IO-bound, not compute-bound",
    "Standard approach: gather tokens \u2192 pad \u2192 GEMM \u2192 scatter back",
]
add_bullet_list(sl, Inches(0.9), Inches(1.5), Inches(11), Inches(3.0),
                items, font_size=20, color=WHITE)

flow = ["Gather", "\u2192", "Pad", "\u2192", "GEMM", "\u2192", "Scatter"]
for i, item in enumerate(flow):
    x = Inches(1.0 + i * 1.6)
    if item == "\u2192":
        add_textbox(sl, x, Inches(4.7), Inches(0.5), Inches(0.5),
                    "\u2192", font_size=28, color=AMBER, bold=True, alignment=PP_ALIGN.CENTER)
    else:
        add_rounded_rect(sl, x, Inches(4.6), Inches(1.4), Inches(0.7),
                         RGBColor(0xF0, 0xF4, 0xF8), item, 16, WHITE, AMBER)

add_textbox(sl, Inches(1.5), Inches(5.8), Inches(10.3), Inches(0.6),
            '"The IO cost scales linearly with expert granularity"',
            font_size=20, color=AMBER, bold=True, alignment=PP_ALIGN.CENTER)

slide_number(sl, slide_num)


# ===== SLIDE 25: SonicMoE Key Innovations =====
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
slide_title(sl, "SonicMoE\u2019s Key Innovations")

innovations = [
    ("Gather Fusion", "Token gather fused with GEMM prologue (no separate kernel)", GREEN),
    ("Ping-Pong Scheduling", "Overlap epilogue IO with next tile\u2019s MMA", BLUE),
    ("Epilogue Fusion", "SwiGLU + dS + A\u2032 computed in GEMM epilogue", AMBER),
]

for i, (name, desc, accent) in enumerate(innovations):
    y = Inches(1.6 + i * 1.5)
    num_box = add_rounded_rect(sl, Inches(0.9), y, Inches(0.6), Inches(0.6),
                               accent, str(i + 1), 20, WHITE)
    add_textbox(sl, Inches(1.8), y, Inches(4.0), Inches(0.5),
                name, font_size=22, color=accent, bold=True)
    add_textbox(sl, Inches(1.8), y + Inches(0.5), Inches(10), Inches(0.5),
                desc, font_size=17, color=WHITE)

add_rounded_rect(sl, Inches(3.0), Inches(5.7), Inches(7.3), Inches(0.8),
                 RGBColor(0xF0, 0xF4, 0xF8),
                 "Result: IO is hidden behind compute",
                 font_size=22, text_color=GREEN, border_color=GREEN)

slide_number(sl, slide_num)


# ===== SLIDE 26: SonicMoE Profiler Results =====
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
slide_title(sl, "SonicMoE \u2014 Profiler Results")

add_textbox(sl, Inches(0.9), Inches(1.5), Inches(11), Inches(0.5),
            "Configuration: E=128, K=8, d=1536, n=256 (7B MoE)",
            font_size=18, color=LIGHT_GRAY)

add_rounded_rect(sl, Inches(0.9), Inches(2.3), Inches(4.0), Inches(0.8),
                 RGBColor(0xF0, 0xF4, 0xF8), "35 TFLOPS on H200",
                 font_size=24, text_color=AMBER, border_color=AMBER)

items = [
    "quackgemm_actGemmGatedSm90  (up projection + SwiGLU fused)",
    "quackgemm_default_epiGemmDefaultSm90  (down projection)",
]
tf = add_textbox(sl, Inches(0.9), Inches(3.6), Inches(11), Inches(0.4),
                 "Key Kernels:", font_size=18, color=BLUE, bold=True)
add_bullet_list(sl, Inches(0.9), Inches(4.1), Inches(11), Inches(1.5),
                items, font_size=17, color=WHITE)

add_textbox(sl, Inches(1.5), Inches(5.8), Inches(10.3), Inches(0.6),
            '"Two GEMM kernels with fused activation \u2014 that\u2019s the entire MoE forward"',
            font_size=19, color=AMBER, bold=True, alignment=PP_ALIGN.CENTER)

slide_number(sl, slide_num)


# ===== SLIDE 27: SonicMoE Perfetto Trace =====
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
slide_title(sl, "SonicMoE \u2014 Perfetto Trace")

add_demo_placeholder(sl, "[LIVE DEMO \u2014 Perfetto UI]\n\nsonicmoe_forward annotation, two GEMM kernels, routing overhead")

slide_number(sl, slide_num)


# ===== SLIDE 28: SonicMoE — NCU Details (NEW) =====
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
slide_title(sl, "SonicMoE \u2014 NCU Details")

add_screenshot_placeholder(sl, "[SCREENSHOT: NCU SonicMoE kernel]",
                           Inches(1.5), Inches(1.5), Inches(10.3), Inches(4.5))

add_textbox(sl, Inches(1.0), Inches(6.3), Inches(11.3), Inches(0.5),
            "quackgemm kernel: occupancy, memory throughput, warp stall breakdown",
            font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

slide_number(sl, slide_num)


# ===== SLIDE 29: The Production Picture =====
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
slide_title(sl, "The Production Picture")

items = [
    "DGX Spark profiling: DeepGEMM accounts for 61.9% of GPU time in DeepSeek-V4-Flash inference",
    "SonicMoE beats DeepGEMM by 43% on forward pass for fine-grained MoEs (from their paper)",
]
add_bullet_list(sl, Inches(0.9), Inches(1.5), Inches(11), Inches(2.5),
                items, font_size=20, color=WHITE)

add_rounded_rect(sl, Inches(2.5), Inches(4.2), Inches(8.3), Inches(1.0),
                 RGBColor(0xF0, 0xF4, 0xF8),
                 "These GEMM kernels ARE the bottleneck in production",
                 font_size=24, text_color=AMBER, border_color=AMBER)

slide_number(sl, slide_num)


# ===== SLIDE 30: Summary Table =====
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
slide_title(sl, "Summary")

data = [
    ["Kernel", "What", "Speedup", "Key NCU Insight"],
    ["DeepGEMM FP8", "JIT FP8 GEMM", "3.09x vs bf16",
     "67% compute, TMA-only data path"],
    ["FlashAttention-3", "Hopper attention", "2.0\u20132.6x vs FA-2",
     "WGMMA + TMA + warp specialization"],
    ["SonicMoE", "IO-aware MoE", "35 TFLOPS",
     "Gather fusion + ping-pong scheduling"],
]

add_table(sl, Inches(0.8), Inches(1.5), Inches(11.7), Inches(2.8),
          len(data), 4, data,
          col_widths=[Inches(2.8), Inches(2.5), Inches(2.5), Inches(3.9)])

slide_number(sl, slide_num)


# ===== SLIDE 31: Takeaways =====
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
slide_title(sl, "Takeaways")

takeaways = [
    ("torch.profiler tells you WHAT happened", "Open the trace in Perfetto", BLUE),
    ("Nsight Compute tells you WHY", "Roofline, occupancy, memory throughput", GREEN),
    ("Kernel names are fingerprints", "sm90, fp8, 1d2d, device_kernel, quackgemm", AMBER),
    ("Profile first. Then optimize. Then profile again.", "", WHITE),
]

for i, (main, sub, accent) in enumerate(takeaways):
    y = Inches(1.5 + i * 1.35)
    num_box = add_rounded_rect(sl, Inches(0.9), y, Inches(0.6), Inches(0.6),
                               accent, str(i + 1), 20, WHITE)
    add_textbox(sl, Inches(1.8), y, Inches(10.5), Inches(0.5),
                main, font_size=22, color=accent, bold=True)
    if sub:
        add_textbox(sl, Inches(1.8), y + Inches(0.5), Inches(10.5), Inches(0.4),
                    sub, font_size=16, color=LIGHT_GRAY)

slide_number(sl, slide_num)


# ===== SLIDE 32: Resources =====
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
slide_title(sl, "Resources")

resources = [
    ("\u25B6  Workshop Repo", "github.com/XAheli/torch-profiler", BLUE),
    ("\u25B6  DeepGEMM", "github.com/deepseek-ai/DeepGEMM", BLUE),
    ("\u25B6  FlashAttention-3", "github.com/Dao-AILab/flash-attention", GREEN),
    ("\u25B6  SonicMoE", "github.com/Dao-AILab/sonic-moe", AMBER),
    ("\u25B6  Perfetto UI", "ui.perfetto.dev", LIGHT_GRAY),
]

for i, (label, url, accent) in enumerate(resources):
    y = Inches(1.5 + i * 0.9)
    add_textbox(sl, Inches(1.2), y, Inches(4.0), Inches(0.5),
                label, font_size=20, color=accent, bold=True)
    add_textbox(sl, Inches(5.5), y, Inches(7.0), Inches(0.5),
                url, font_size=18, color=WHITE)

add_accent_bar(sl, Inches(1.2), Inches(6.2), Inches(11), Inches(0.03), BLUE)
add_textbox(sl, Inches(1.2), Inches(6.4), Inches(11), Inches(0.5),
            "Aheli Poddar  |  Red Hat PyTorch Team", font_size=18, color=LIGHT_GRAY,
            alignment=PP_ALIGN.CENTER)

slide_number(sl, slide_num)


# ===== SLIDE 33: Thank You =====
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)

add_accent_bar(sl, Inches(0), Inches(0), SLIDE_W, Inches(0.06), BLUE)

add_textbox(sl, Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.0),
            "How Fast Is Your Kernel, Really?",
            font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_accent_bar(sl, Inches(5.5), Inches(3.2), Inches(2.3), Inches(0.04), BLUE)

add_textbox(sl, Inches(1.5), Inches(3.6), Inches(10.3), Inches(0.6),
            "Thank You!", font_size=36, color=BLUE, bold=True,
            alignment=PP_ALIGN.CENTER)

add_textbox(sl, Inches(1.5), Inches(4.6), Inches(10.3), Inches(0.5),
            "Aheli Poddar  |  Red Hat PyTorch Team", font_size=22, color=LIGHT_GRAY,
            alignment=PP_ALIGN.CENTER)

add_textbox(sl, Inches(1.5), Inches(5.5), Inches(10.3), Inches(0.5),
            "Questions?", font_size=28, color=AMBER, bold=True,
            alignment=PP_ALIGN.CENTER)

add_accent_bar(sl, Inches(0), Inches(7.44), SLIDE_W, Inches(0.06), BLUE)
slide_number(sl, slide_num)


# ---------------------------------------------------------------------------
# Save
# ---------------------------------------------------------------------------
prs.save(OUT)
print(f"Saved {slide_num} slides to {OUT}")
